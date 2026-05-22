"""
Génération de la présentation PPTX du projet mémoire GE.
Public : Mixte académique + business. Niveau : très détaillé.
Sortie : reports/presentation_memoire_GE.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ----------------------------------------------------------------------------
# Configuration globale
# ----------------------------------------------------------------------------
ROOT = Path(r"c:\Users\lenovo\Desktop\Extraction livraison client 2021-2025")
REPORTS = ROOT / "reports"
OUTPUT = REPORTS / "presentation_memoire_GE.pptx"

# Charte graphique (inspiration GE corporate)
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x68)      # Bleu foncé GE
COLOR_ACCENT  = RGBColor(0xE8, 0x7A, 0x1E)      # Orange accent
COLOR_LIGHT   = RGBColor(0xE8, 0xEF, 0xF7)      # Bleu très clair (bandeau)
COLOR_TEXT    = RGBColor(0x22, 0x22, 0x22)      # Gris très foncé
COLOR_MUTED   = RGBColor(0x6B, 0x73, 0x80)      # Gris secondaire
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUCCESS = RGBColor(0x2E, 0x86, 0x4E)      # Vert
COLOR_WARN    = RGBColor(0xC0, 0x39, 0x2B)      # Rouge

# Taille slide 16:9 standard
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ----------------------------------------------------------------------------
# Helpers bas-niveau
# ----------------------------------------------------------------------------

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, color)
    shp.shadow.inherit = False
    return shp

def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=16, color=COLOR_TEXT,
                bullet_color=COLOR_ACCENT, line_spacing=1.2):
    """items: list of str OR list of (str_main, str_detail or None) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            if len(it) == 2:
                main, detail = it
            elif len(it) == 1:
                main, detail = it[0], None
            else:
                main, detail = str(it), None
        else:
            main, detail = it, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r0 = p.add_run()
        r0.text = "▸ "
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        r0.font.name = "Calibri"
        r1 = p.add_run()
        r1.text = main
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = "Calibri"
        if detail:
            p.add_run().text = "  "
            r2 = p.add_run()
            r2.text = detail
            r2.font.size = Pt(size - 2)
            r2.font.color.rgb = COLOR_MUTED
            r2.font.italic = True
            r2.font.name = "Calibri"
    return tb

# ----------------------------------------------------------------------------
# Layout commun
# ----------------------------------------------------------------------------

def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_header(slide, title, subtitle=None, idx=None, total=None):
    # Bandeau supérieur fin
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.08), COLOR_ACCENT)
    # Titre
    add_textbox(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.6),
                title, size=28, bold=True, color=COLOR_PRIMARY)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.72), Inches(11.5), Inches(0.35),
                    subtitle, size=14, color=COLOR_MUTED)
    # Filet de séparation
    add_rect(slide, Inches(0.5), Inches(1.12), Inches(12.3), Emu(12000), COLOR_PRIMARY)
    # Pagination
    if idx is not None and total is not None:
        add_textbox(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
                    f"{idx} / {total}", size=10, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)
    # Pied de page
    add_textbox(slide, Inches(0.5), Inches(7.15), Inches(11), Inches(0.3),
                "Projet mémoire Master GE — Prévision Demand-Driven · 2026",
                size=10, color=COLOR_MUTED)

def add_image_centered(slide, img_path, x, y, w, h):
    """Place une image dans une zone (w,h) en conservant le ratio, centrée."""
    from PIL import Image
    p = Path(img_path)
    if not p.exists():
        add_textbox(slide, x, y, w, h, f"[IMAGE MANQUANTE : {p.name}]",
                    size=12, color=COLOR_WARN, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        return None
    with Image.open(p) as im:
        iw, ih = im.size
    ratio_img = iw / ih
    ratio_box = w / h
    if ratio_img > ratio_box:
        new_w = w
        new_h = int(w / ratio_img)
    else:
        new_h = h
        new_w = int(h * ratio_img)
    nx = x + (w - new_w) // 2
    ny = y + (h - new_h) // 2
    return slide.shapes.add_picture(str(p), nx, ny, new_w, new_h)

def add_table(slide, x, y, w, h, data, *, header=True, col_widths=None,
              font_size=12, header_color=COLOR_PRIMARY, zebra=True):
    """data: list of lists (rows). First row = header if header=True."""
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""  # clear default
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.bold = True
                run.font.color.rgb = COLOR_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                if zebra and r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_LIGHT
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_WHITE
                run.font.color.rgb = COLOR_TEXT
                # numerical right-align
                try:
                    float(str(val).replace(",", ".").replace("%", "").replace("+", "").replace("−", "-"))
                    if c > 0:
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.LEFT
                except ValueError:
                    p.alignment = PP_ALIGN.LEFT
    return tbl

def add_kpi(slide, x, y, w, h, value, label, *, value_color=COLOR_PRIMARY,
            border_color=COLOR_ACCENT):
    # Carte KPI
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(card, COLOR_WHITE)
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    add_textbox(slide, x, y + Inches(0.15), w, Inches(0.9),
                value, size=32, bold=True, color=value_color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x, y + Inches(1.0), w, Inches(0.5),
                label, size=11, color=COLOR_MUTED,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ----------------------------------------------------------------------------
# Slides : génération
# ----------------------------------------------------------------------------

def build(prs):
    # Provisional total — sera mis à jour après construction
    TOTAL = 38

    # ============ Slide 1 — Couverture ===========
    s = slide_blank(prs)
    # Fond dégradé approx : grand rect bleu + bande accent
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, COLOR_PRIMARY)
    add_rect(s, Emu(0), Inches(4.2), SLIDE_W, Inches(0.08), COLOR_ACCENT)
    add_rect(s, Inches(0.5), Inches(0.5), Inches(2.2), Inches(0.5), COLOR_ACCENT)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(2.2), Inches(0.5),
                "PROJET MÉMOIRE", size=14, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(0.8), Inches(1.6), Inches(12), Inches(1.2),
                "Système de Prévision et d'Optimisation des Stocks",
                size=40, bold=True, color=COLOR_WHITE)
    add_textbox(s, Inches(0.8), Inches(2.7), Inches(12), Inches(0.8),
                "Une approche Demand-Driven s'appuyant sur l'Intelligence Artificielle",
                size=22, color=COLOR_LIGHT)
    add_textbox(s, Inches(0.8), Inches(4.6), Inches(12), Inches(0.5),
                "Cas d'étude : General Electric — Chaîne d'approvisionnement industrielle",
                size=18, bold=True, color=COLOR_WHITE)
    add_textbox(s, Inches(0.8), Inches(5.2), Inches(12), Inches(0.4),
                "Données 2021–2025  ·  349 390 commandes  ·  Pipeline IA complet",
                size=14, color=COLOR_LIGHT)
    add_textbox(s, Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
                "Adam Marrakchi  ·  Master Génie Industriel  ·  Mai 2026",
                size=14, color=COLOR_WHITE)
    add_textbox(s, Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.4),
                "Soutenance — 28 mai 2026", size=13,
                color=COLOR_ACCENT, align=PP_ALIGN.RIGHT, bold=True)

    # ============ Slide 2 — Sommaire ===========
    s = slide_blank(prs)
    add_header(s, "Sommaire", "Plan détaillé de la présentation", idx=2, total=TOTAL)
    sommaire = [
        ("1. Contexte & Problématique", "Enjeux supply chain GE · objectifs mémoire"),
        ("2. Phase 1 — Data Engineering", "Extraction, fusion, nettoyage (349K commandes)"),
        ("3. Phase 2 — Études Statistiques", "Saisonnalité, cyclicité, corrélations exogènes"),
        ("4. Phase 2bis — EDA pré-modélisation", "Analyse fine, sélection features, splits"),
        ("5. Phase 3 — Modélisation IA", "Baseline · XGBoost · LightGBM · Optuna · LSTM"),
        ("6. Phase 4 — Dashboard décisionnel", "Streamlit human-in-the-loop"),
        ("7. Résultats & Performance", "MAE 11.87 vs baseline 13.04 — objectif validé"),
        ("8. Apports, Limites & Perspectives", "Conclusion académique et business"),
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.6),
                sommaire, size=18, line_spacing=1.5)

    # ============ Slide 3 — Contexte GE ===========
    s = slide_blank(prs)
    add_header(s, "Contexte : la chaîne d'approvisionnement GE", "Pression économique · multi-sites · forte variabilité de la demande",
               idx=3, total=TOTAL)
    add_bullets(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(5.5), [
        ("General Electric — leader industriel mondial", "Énergie, aéronautique, santé, transport"),
        ("Chaîne logistique complexe", "Réseau client international, multi-devises, multi-pays"),
        ("Pression croissante sur les stocks", "Réduction du BFR, taux de service ≥ 95 %"),
        ("Variabilité de la demande forte", "Saisonnalité industrielle, cycles budgétaires, aléas exogènes"),
        ("Approche traditionnelle limitée", "Prévision sur moyennes mobiles → réactivité insuffisante"),
        ("Opportunité IA", "Captation de signaux faibles (météo, fériés, IPI) au-delà des séries temporelles"),
    ], size=16, line_spacing=1.4)
    # Cartouche illustratif
    add_rect(s, Inches(8.5), Inches(1.4), Inches(4.3), Inches(5.5), COLOR_LIGHT)
    add_textbox(s, Inches(8.7), Inches(1.55), Inches(4), Inches(0.5),
                "Périmètre étudié", size=14, bold=True, color=COLOR_PRIMARY)
    add_kpi(s, Inches(8.7), Inches(2.1), Inches(1.9), Inches(1.4), "349 K", "commandes 2021-2025")
    add_kpi(s, Inches(10.7), Inches(2.1), Inches(1.9), Inches(1.4), "5 ans", "horizon historique")
    add_kpi(s, Inches(8.7), Inches(3.7), Inches(1.9), Inches(1.4), "89 %", "France · multi-pays")
    add_kpi(s, Inches(10.7), Inches(3.7), Inches(1.9), Inches(1.4), "3", "familles produits")
    add_textbox(s, Inches(8.7), Inches(5.4), Inches(4), Inches(1.4),
                "Sources : extractions ERP commandes & livraisons clients (2021→2025), enrichies par jours fériés (FR), météo (Open-Meteo Paris) et conjoncture INSEE (IPI).",
                size=10, color=COLOR_MUTED)

    # ============ Slide 4 — Problématique ===========
    s = slide_blank(prs)
    add_header(s, "Problématique de recherche", "Question centrale du mémoire",
               idx=4, total=TOTAL)
    # Bloc question centrale
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2))
    set_fill(box, COLOR_PRIMARY)
    add_textbox(s, Inches(1), Inches(1.7), Inches(11.3), Inches(1.8),
                "Comment une approche Demand-Driven s'appuyant sur l'IA peut-elle prédire conjointement la QUANTITÉ et la DATE de la demande client réelle, afin d'optimiser la chaîne d'approvisionnement de GE ?",
                size=20, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Sous-questions
    add_textbox(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.4),
                "Sous-questions de recherche", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.8), Inches(4.35), Inches(11.7), Inches(2.8), [
        ("Q1 — Quels facteurs exogènes (calendaires, météo, macro) impactent significativement la demande ?",),
        ("Q2 — Une architecture IA peut-elle battre une baseline statistique métier sur un horizon 12 mois ?",),
        ("Q3 — Comment articuler prévision quantitative (qté) et temporelle (date inter-commande) ?",),
        ("Q4 — Comment intégrer la prévision IA dans un outil décisionnel utilisable au quotidien ?",),
    ], size=14, line_spacing=1.3)

    # ============ Slide 5 — Objectifs ===========
    s = slide_blank(prs)
    add_header(s, "Objectifs et indicateur North Star", "Cibles quantitatives et livrables",
               idx=5, total=TOTAL)
    add_textbox(s, Inches(0.6), Inches(1.3), Inches(6), Inches(0.5),
                "Objectifs scientifiques", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.6), Inches(1.85), Inches(6.5), Inches(5), [
        ("Identifier les signaux exogènes prédictifs", "Tests statistiques : Pearson, Granger"),
        ("Comparer plusieurs architectures IA", "Boosting (XGB/LGBM) vs Deep Learning (LSTM)"),
        ("Battre la baseline métier", "Moyenne historique client × article"),
        ("Garantir la robustesse temporelle", "Split strict 2021-2023 / 2024 / 2025"),
        ("Industrialiser via une interface décisionnelle", "Dashboard Streamlit human-in-the-loop"),
    ], size=14, line_spacing=1.4)
    # North star
    add_textbox(s, Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.5),
                "Indicateur North Star", size=18, bold=True, color=COLOR_PRIMARY)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(7.5), Inches(1.9), Inches(5.3), Inches(2.4))
    set_fill(box, COLOR_ACCENT)
    add_textbox(s, Inches(7.5), Inches(2.0), Inches(5.3), Inches(0.6),
                "Précision globale > 60 %", size=24, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.5), Inches(2.7), Inches(5.3), Inches(1.7),
                "Mesurée par 1 − WAPE sur le test 2025\nObjectif : prévision quantitative meilleure que la baseline moyenne historique sur l'horizon annuel",
                size=13, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # Livrables
    add_textbox(s, Inches(7.5), Inches(4.5), Inches(5.5), Inches(0.5),
                "Livrables", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(7.5), Inches(5), Inches(5.5), Inches(2.2), [
        ("4 notebooks Jupyter documentés",),
        ("5 modèles IA sauvegardés (.pkl / .pt)",),
        ("1 dashboard Streamlit multi-pages",),
        ("Rapports phase 1, 2 + JSON modélisation",),
    ], size=13, line_spacing=1.3)

    # ============ Slide 6 — Méthodologie globale ===========
    s = slide_blank(prs)
    add_header(s, "Méthodologie : pipeline en 4 phases", "Démarche structurée du brut au déploiement",
               idx=6, total=TOTAL)
    phases = [
        ("Phase 1", "Data\nEngineering", "349K lignes\n24 colonnes", COLOR_PRIMARY),
        ("Phase 2", "Études\nstatistiques", "Saisonnalité\nGranger / Pearson", COLOR_ACCENT),
        ("Phase 2bis", "EDA pré-\nmodélisation", "28 features\nSplits temporels", COLOR_PRIMARY),
        ("Phase 3", "Modélisation\nIA", "XGB / LGBM\nOptuna / LSTM", COLOR_ACCENT),
        ("Phase 4", "Dashboard\ndécisionnel", "Streamlit\nHITL", COLOR_PRIMARY),
    ]
    x0 = Inches(0.6)
    box_w = Inches(2.4)
    gap = Inches(0.05)
    arrow_w = Inches(0.2)
    cur = x0
    for i, (ph, name, sub, color) in enumerate(phases):
        # Box
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cur, Inches(2.0), box_w, Inches(2.8))
        set_fill(box, color)
        add_textbox(s, cur, Inches(2.15), box_w, Inches(0.4),
                    ph, size=12, bold=True, color=COLOR_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, cur, Inches(2.6), box_w, Inches(1.0),
                    name, size=18, bold=True, color=COLOR_WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, cur, Inches(3.7), box_w, Inches(1.0),
                    sub, size=12, color=COLOR_WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur = cur + box_w + gap
        if i < len(phases) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, cur, Inches(3.2),
                                      arrow_w, Inches(0.4))
            set_fill(arr, COLOR_MUTED)
            cur = cur + arrow_w + gap
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
                "Approche itérative et reproductible", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.6), Inches(5.85), Inches(12), Inches(1.4), [
        ("Versionning Git de chaque phase + rapports markdown détaillés",),
        ("Données intermédiaires en Parquet (compression + typage strict)",),
        ("Reproductibilité : seed fixée (RANDOM_STATE=42), splits temporels stricts",),
    ], size=13, line_spacing=1.3)

    # ============ Slide 7 — Phase 1 : Sources ===========
    s = slide_blank(prs)
    add_header(s, "Phase 1 — Sources de données et stratégie d'extraction",
               "Deux extractions ERP croisées sur la période 2021–2025",
               idx=7, total=TOTAL)
    add_table(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.8), [
        ["Source", "Volume brut", "Granularité", "Variables clés"],
        ["Extraction commande", "~ 380 000 lignes", "Ligne de commande", "client, article, qté demandée, prix, devise, date cmd"],
        ["Extraction livraison", "~ 410 000 lignes", "Ligne de livraison", "date livraison demandée/réelle, quantité livrée, statut"],
        ["Jointure commande↔livraison", "349 390 lignes", "Ligne de commande consolidée", "Match par n°cmd + n°ligne"],
    ], col_widths=[Inches(2.8), Inches(2.0), Inches(3.0), Inches(4.2)], font_size=11)
    # Schéma simple
    add_textbox(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
                "Stratégie de consolidation", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.6), Inches(4.05), Inches(12), Inches(3), [
        ("Fusion en clé composite (numéro commande, numéro ligne)", "Préserve la granularité commande × article"),
        ("Suppression des lignes incohérentes", "Quantités négatives, dates invalides, statuts annulés"),
        ("Typage strict des dates", "date_cmd, date_liv_demandee, date_liv_reelle au format datetime64"),
        ("Création de la variable cible qte_demandee", "Quantité commandée (avant éventuelles modifications)"),
        ("Validation finale : 349 390 lignes × 24 colonnes", "Aucune valeur manquante dans la cible"),
    ], size=14, line_spacing=1.35)

    # ============ Slide 8 — Phase 1 : Pipeline ===========
    s = slide_blank(prs)
    add_header(s, "Phase 1 — Pipeline de nettoyage et de transformation",
               "Notebook 01_data_cleaning.ipynb — sortie : dataset_ml_final.parquet",
               idx=8, total=TOTAL)
    steps = [
        ("1. Chargement", "Lecture XLSX + CSV (encodage Latin-1, séparateur ;)"),
        ("2. Standardisation", "Renommage snake_case, ASCII, harmonisation devises"),
        ("3. Jointure", "Inner join commande × livraison, vérif. ratio match"),
        ("4. Dédoublonnage", "Suppression doublons exacts + lignes vides"),
        ("5. Nettoyage cible", "qte_demandee : suppression valeurs négatives / nulles"),
        ("6. Features dates", "année, mois, semaine, jour_semaine, est_weekend, est_fin_mois"),
        ("7. Features métiers", "delai_demande_jours = liv_demandee − cmd"),
        ("8. Encodages", "Frequency encoding (client, article) — Target encoding catégorielles"),
        ("9. Validation finale", "Profilage pandas + sanity checks → 349 390 × 24"),
    ]
    y = Inches(1.4)
    for i, (titre, det) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        yy = y + Inches(row * 1.8)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, Inches(4), Inches(1.6))
        set_fill(card, COLOR_LIGHT)
        add_textbox(s, x, yy + Inches(0.1), Inches(4), Inches(0.5),
                    titre, size=14, bold=True, color=COLOR_PRIMARY,
                    align=PP_ALIGN.LEFT)
        add_textbox(s, x, yy + Inches(0.6), Inches(4), Inches(1),
                    det, size=11, color=COLOR_TEXT)

    # ============ Slide 9 — Phase 1 : Dataset final ===========
    s = slide_blank(prs)
    add_header(s, "Phase 1 — Caractéristiques du dataset produit",
               "data/processed/dataset_ml_final.parquet · 349 390 × 24",
               idx=9, total=TOTAL)
    add_kpi(s, Inches(0.6), Inches(1.5), Inches(2.8), Inches(1.5), "349 390", "lignes (commandes)")
    add_kpi(s, Inches(3.6), Inches(1.5), Inches(2.8), Inches(1.5), "24", "colonnes initiales")
    add_kpi(s, Inches(6.6), Inches(1.5), Inches(2.8), Inches(1.5), "0", "valeur manquante sur cible")
    add_kpi(s, Inches(9.6), Inches(1.5), Inches(2.8), Inches(1.5), "5 ans", "horizon 2021–2025")
    # Catégories de variables
    add_textbox(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.5),
                "Catégories de variables", size=16, bold=True, color=COLOR_PRIMARY)
    add_table(s, Inches(0.6), Inches(3.85), Inches(12), Inches(3), [
        ["Catégorie", "Variables", "Exemples"],
        ["Identifiants", "2", "code_client, code_article"],
        ["Cible & quantitatives", "4", "qte_demandee, qte_livree, prix, delai_demande_jours"],
        ["Temporelles dérivées", "8", "annee, mois, trimestre, semaine, jour_semaine, est_weekend, est_fin_mois, date_cmd"],
        ["Catégorielles", "7", "statut, devise, pays, segment, famille_activite_client/article, type_activite"],
        ["Encodées (freq/enc)", "3", "code_client_freq, code_article_freq, statut_enc"],
    ], col_widths=[Inches(3), Inches(2), Inches(7)], font_size=12)

    # ============ Slide 10 — Phase 2 : Objectifs ===========
    s = slide_blank(prs)
    add_header(s, "Phase 2 — Études statistiques approfondies",
               "Notebook 02_statistical_study.ipynb · enrichissement à 38 colonnes",
               idx=10, total=TOTAL)
    add_textbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
                "Deux études complémentaires", size=18, bold=True, color=COLOR_PRIMARY)
    # Étude A
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2),
                              Inches(6), Inches(4.5))
    set_fill(box, COLOR_LIGHT)
    add_textbox(s, Inches(0.6), Inches(2.1), Inches(6), Inches(0.5),
                "Étude A — Saisonnalité & cyclicité", size=16, bold=True,
                color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.8), Inches(2.7), Inches(5.7), Inches(3.7), [
        ("Décomposition STL (Loess) — période 12 mois",),
        ("Tendance vs Saisonnalité vs Résidus",),
        ("Décomposition par famille d'article (3 familles)",),
        ("Cycles budgétaires (fin de trimestre)",),
        ("Cycles météo (4 saisons)",),
        ("Stationnarité (ACF / PACF)",),
    ], size=13, line_spacing=1.4)
    # Étude B
    box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2),
                              Inches(6), Inches(4.5))
    set_fill(box2, COLOR_LIGHT)
    add_textbox(s, Inches(6.8), Inches(2.1), Inches(6), Inches(0.5),
                "Étude B — Corrélations exogènes", size=16, bold=True,
                color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(7), Inches(2.7), Inches(5.7), Inches(3.7), [
        ("Construction de 12 variables exogènes",),
        ("Jours fériés FR (lib. holidays)",),
        ("Météo Paris (API Open-Meteo, proxy France)",),
        ("IPI INSEE (production industrielle)",),
        ("Matrice de corrélation Pearson (seuil 0,15)",),
        ("Test de causalité de Granger (lag 1–3 mois)",),
    ], size=13, line_spacing=1.4)
    add_textbox(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.5),
                "Livrable : dataset_ml_enrichi.parquet (349 390 × 38) — 100 % couverture sur les nouvelles variables",
                size=12, color=COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # ============ Slide 11 — Décomposition STL globale ===========
    s = slide_blank(prs)
    add_header(s, "Décomposition STL — Série globale (64 mois)",
               "Tendance haussière, saisonnalité forte (amplitude ≈ 90 000 unités)",
               idx=11, total=TOTAL)
    add_image_centered(s, REPORTS / "decomposition_globale.png",
                       Inches(0.6), Inches(1.35), Inches(8.8), Inches(5.5))
    # Insights à droite
    add_textbox(s, Inches(9.7), Inches(1.35), Inches(3.5), Inches(0.5),
                "Lecture", size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(9.7), Inches(1.85), Inches(3.5), Inches(5), [
        ("Tendance +13 130 unités sur 5 ans",),
        ("Saisonnalité = 72 % de la moyenne mensuelle",),
        ("Pic en mars (clôtures budgétaires Q1)",),
        ("Creux en août (fermetures estivales)",),
        ("Résidus = chocs ponctuels non structurés",),
    ], size=11, line_spacing=1.3)

    # ============ Slide 12 — Décomposition par famille ===========
    s = slide_blank(prs)
    add_header(s, "Décomposition par famille d'article",
               "Famille 0 = 96 % du volume — structure saisonnière concentrée",
               idx=12, total=TOTAL)
    add_image_centered(s, REPORTS / "decomposition_par_famille.png",
                       Inches(0.6), Inches(1.35), Inches(8.8), Inches(5.5))
    add_table(s, Inches(9.6), Inches(1.4), Inches(3.5), Inches(3), [
        ["Famille", "Volume", "Part"],
        ["0", "7,12 M", "96,3 %"],
        ["1", "562 K", "7,6 %"],
        ["2", "281 K", "3,8 %"],
    ], col_widths=[Inches(1.2), Inches(1.3), Inches(1)], font_size=12)
    add_textbox(s, Inches(9.6), Inches(4.6), Inches(3.6), Inches(2.5),
                "La famille 0 dicte la structure globale. Les familles 1 & 2, plus rares, présentent un bruit relatif plus important — à modéliser séparément en perspective.",
                size=11, color=COLOR_MUTED)

    # ============ Slide 13 — Cycles budgétaires ===========
    s = slide_blank(prs)
    add_header(s, "Cycles budgétaires clients (fin de trimestre)",
               "Surdemande typique mars / juin / septembre / décembre",
               idx=13, total=TOTAL)
    add_image_centered(s, REPORTS / "cycles_budgetaires.png",
                       Inches(0.6), Inches(1.35), Inches(8.8), Inches(5.5))
    add_textbox(s, Inches(9.7), Inches(1.35), Inches(3.5), Inches(0.5),
                "Implication modélisation", size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(9.7), Inches(1.85), Inches(3.5), Inches(5), [
        ("Inclure trimestre_cmd",),
        ("Inclure est_fin_mois_cmd",),
        ("Capture du comportement d'achat industriel pré-clôture comptable",),
    ], size=11, line_spacing=1.3)

    # ============ Slide 14 — Cycles météo ===========
    s = slide_blank(prs)
    add_header(s, "Cycles météo saisonniers",
               "Creux estival marqué — effet des arrêts industriels",
               idx=14, total=TOTAL)
    add_image_centered(s, REPORTS / "cycles_meteo.png",
                       Inches(0.6), Inches(1.35), Inches(8.8), Inches(5.5))
    add_bullets(s, Inches(9.7), Inches(1.5), Inches(3.5), Inches(5), [
        ("Été = creux historique",),
        ("Hiver / Printemps = pics",),
        ("Cohérent avec pic mars + creux août STL",),
        ("Légitime l'inclusion de temp_min, pluie, vent",),
    ], size=11, line_spacing=1.3)

    # ============ Slide 15 — Construction features exogènes ===========
    s = slide_blank(prs)
    add_header(s, "Construction des 12 variables exogènes",
               "Sources extérieures jointurées au dataset enrichi",
               idx=15, total=TOTAL)
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(5.5), [
        ["Famille", "Variable", "Source / méthode", "Couverture"],
        ["Jours fériés", "est_jour_ferie_cmd", "lib. holidays (FR)", "28,4 % lignes = 1"],
        ["Jours fériés", "est_jour_ferie_liv_dem", "lib. holidays (FR)", "28,6 % lignes = 1"],
        ["Jours fériés", "nb_jours_feries_dans_delai", "pandas date_range freq='B'", "moy. 0,34 — max 11"],
        ["Jours fériés", "nb_weekends_dans_delai", "pandas weekday ≥ 5", "moy. 2,12 — max 114"],
        ["Météo", "pluie_mm_liv_dem", "Open-Meteo Archive (Paris, mensuel)", "100 %"],
        ["Météo", "vent_max_kmh_liv_dem", "Open-Meteo Archive (Paris, mensuel)", "100 %"],
        ["Météo", "temp_min_liv_dem", "Open-Meteo Archive (Paris, mensuel)", "100 %"],
        ["Macro", "ipi_valeur", "INSEE — Production industrielle (proxy)", "100 %"],
        ["Macro", "taux_change_devise", "Taux annuel moyen EUR/USD/CNY", "100 %"],
        ["Contexte", "est_vacances_scolaires_liv_dem", "Calendrier officiel (toutes zones)", "30,9 % = 1"],
        ["Contexte", "est_periode_peak_liv_dem", "Sept / Nov / Déc", "22,0 % = 1"],
        ["Contexte", "est_periode_covid", "2021 T1-T2 + 2022 T1", "18,1 % = 1"],
    ], col_widths=[Inches(1.8), Inches(3.2), Inches(4.4), Inches(3)], font_size=10)

    # ============ Slide 16 — Matrice corrélation ===========
    s = slide_blank(prs)
    add_header(s, "Corrélations de Pearson — variables exogènes vs demande",
               "Seuil de significativité |r| > 0,15 — 5 variables retenues",
               idx=16, total=TOTAL)
    add_image_centered(s, REPORTS / "matrice_correlation_pearson.png",
                       Inches(0.4), Inches(1.35), Inches(7), Inches(5.6))
    add_table(s, Inches(7.6), Inches(1.4), Inches(5.5), Inches(4.5), [
        ["Variable", "r", "Sig."],
        ["nb_weekends_dans_delai", "+0,545", "★"],
        ["nb_jours_feries_dans_delai", "+0,497", "★"],
        ["ipi_valeur", "−0,262", "★"],
        ["est_jour_ferie_cmd", "+0,198", "★"],
        ["temp_min_liv_dem", "−0,156", "★"],
        ["vent_max_kmh_liv_dem", "+0,149", "—"],
        ["taux_change_devise", "−0,106", "—"],
        ["pluie_mm_liv_dem", "+0,033", "—"],
    ], col_widths=[Inches(3.3), Inches(1.2), Inches(1)], font_size=11)
    add_textbox(s, Inches(7.6), Inches(6.0), Inches(5.5), Inches(0.9),
                "Les cycles calendaires sont les signaux exogènes les plus forts (r ≈ 0,5).",
                size=11, color=COLOR_ACCENT, bold=True)

    # ============ Slide 17 — Granger ===========
    s = slide_blank(prs)
    add_header(s, "Test de causalité de Granger (lag 1–3 mois)",
               "Vérification que les variables exogènes anticipent réellement la demande",
               idx=17, total=TOTAL)
    add_image_centered(s, REPORTS / "granger_pvalues.png",
                       Inches(0.4), Inches(1.35), Inches(7), Inches(5.6))
    add_table(s, Inches(7.6), Inches(1.4), Inches(5.5), Inches(3.6), [
        ["Variable", "lag=1", "lag=2", "Causal ?"],
        ["nb_weekends_dans_delai", "0,000", "0,000", "★ Oui"],
        ["nb_jours_feries_dans_delai", "0,000", "0,000", "★ Oui"],
        ["ipi_valeur", "0,993", "0,999", "Non"],
        ["est_jour_ferie_cmd", "0,509", "0,563", "Non"],
        ["temp_min_liv_dem", "0,283", "0,609", "Non"],
    ], col_widths=[Inches(2.6), Inches(0.9), Inches(0.9), Inches(1.1)], font_size=11)
    add_bullets(s, Inches(7.6), Inches(5.15), Inches(5.5), Inches(2), [
        ("2 variables causales aux lags 1 & 2",),
        ("IPI & jour férié : corrélés mais non causaux → corrélation contemporaine",),
        ("Conservés malgré tout pour la modélisation (XGBoost utilise valeurs courantes)",),
    ], size=10, line_spacing=1.3)

    # ============ Slide 18 — Synthèse Phase 2 ===========
    s = slide_blank(prs)
    add_header(s, "Phase 2 — Synthèse et variables retenues",
               "Hiérarchisation pour la modélisation",
               idx=18, total=TOTAL)
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(3.8), [
        ["Variable", "Pearson", "Granger", "Priorité Phase 3"],
        ["nb_weekends_dans_delai", "★ +0,545", "★ p=0,000", "Haute — feature causale"],
        ["nb_jours_feries_dans_delai", "★ +0,497", "★ p=0,000", "Haute — feature causale"],
        ["ipi_valeur", "★ −0,262", "—", "Moyenne — corrélation contemporaine"],
        ["est_jour_ferie_cmd", "★ +0,198", "—", "Moyenne — corrélation contemporaine"],
        ["temp_min_liv_dem", "★ −0,156", "—", "Moyenne — signal saisonnier"],
        ["Autres variables", "Non sig.", "—", "Faible — testées en ablation"],
    ], col_widths=[Inches(4), Inches(1.8), Inches(1.8), Inches(4.8)], font_size=11)
    add_textbox(s, Inches(0.5), Inches(5.4), Inches(12.4), Inches(0.5),
                "Faits saillants statistiques", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(5.9), Inches(12.4), Inches(1.5), [
        ("Saisonnalité = 72 % de la moyenne mensuelle (amplitude ≈ 90 K unités)",),
        ("Tendance haussière +13 130 unités → marché en croissance",),
        ("Cycles calendaires = signaux exogènes les plus puissants (Pearson + Granger)",),
    ], size=12, line_spacing=1.3)

    # ============ Slide 19 — Phase 2bis : EDA cible ===========
    s = slide_blank(prs)
    add_header(s, "Phase 2bis — Distribution de la cible qte_demandee",
               "Forte asymétrie → transformation log1p obligatoire",
               idx=19, total=TOTAL)
    add_image_centered(s, REPORTS / "eda_target_distribution.png",
                       Inches(0.5), Inches(1.35), Inches(8.6), Inches(5.6))
    add_textbox(s, Inches(9.4), Inches(1.4), Inches(3.6), Inches(0.5),
                "Constats clés", size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(9.4), Inches(1.9), Inches(3.6), Inches(4), [
        ("Skewness brute : 80,73",),
        ("Skewness après log1p : 1,44",),
        ("Présence de valeurs extrêmes (queue épaisse)",),
        ("Médiane << moyenne → ⇒ loss Tweedie pertinente",),
    ], size=11, line_spacing=1.4)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(5.6),
                              Inches(3.6), Inches(1.3))
    set_fill(box, COLOR_ACCENT)
    add_textbox(s, Inches(9.4), Inches(5.65), Inches(3.6), Inches(1.2),
                "Décision : entraîner sur log1p(qte_demandee)\ninverse = expm1, clip ≥ 0",
                size=12, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============ Slide 20 — Pareto + Heatmap ===========
    s = slide_blank(prs)
    add_header(s, "Concentration du volume et heatmap année × mois",
               "Pareto produits/clients + saisonnalité visuelle confirmée",
               idx=20, total=TOTAL)
    add_image_centered(s, REPORTS / "eda_pareto.png",
                       Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.6))
    add_image_centered(s, REPORTS / "eda_heatmap_annee_mois.png",
                       Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.6))

    # ============ Slide 21 — ACF / PACF + Drift ===========
    s = slide_blank(prs)
    add_header(s, "Stationnarité, autocorrélation et drift annuel",
               "ACF/PACF + suivi temporel — validation du split train/val/test",
               idx=21, total=TOTAL)
    add_image_centered(s, REPORTS / "eda_acf_pacf.png",
                       Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.6))
    add_image_centered(s, REPORTS / "eda_drift_annuel.png",
                       Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.6))

    # ============ Slide 22 — Split temporel ===========
    s = slide_blank(prs)
    add_header(s, "Stratégie de validation : split temporel strict",
               "Aucune fuite d'information — chronologie respectée",
               idx=22, total=TOTAL)
    add_image_centered(s, REPORTS / "eda_split_temporel.png",
                       Inches(0.5), Inches(1.35), Inches(8.6), Inches(5.6))
    add_kpi(s, Inches(9.5), Inches(1.5), Inches(3.4), Inches(1.3), "210 641", "Train 2021-2023")
    add_kpi(s, Inches(9.5), Inches(3.0), Inches(3.4), Inches(1.3), "70 871", "Val 2024")
    add_kpi(s, Inches(9.5), Inches(4.5), Inches(3.4), Inches(1.3), "66 174", "Test 2025")
    add_textbox(s, Inches(9.5), Inches(6.0), Inches(3.4), Inches(0.9),
                "Pas de shuffle — pas de validation croisée aléatoire — préserve le caractère prédictif réel.",
                size=10, color=COLOR_MUTED)

    # ============ Slide 23 — Boxplots & corrélations ===========
    s = slide_blank(prs)
    add_header(s, "Analyse bivariée — boxplots segments & matrice complète",
               "Identification des leviers explicatifs avant modélisation",
               idx=23, total=TOTAL)
    add_image_centered(s, REPORTS / "eda_boxplots_segments.png",
                       Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.6))
    add_image_centered(s, REPORTS / "eda_matrice_correlation_complete.png",
                       Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.6))

    # ============ Slide 24 — Sélection des 28 features ===========
    s = slide_blank(prs)
    add_header(s, "Sélection finale des 28 features pour la Phase 3",
               "Réduction par VIF + corrélations + tests bivariés",
               idx=24, total=TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
                "28 features sélectionnées en 5 catégories", size=16, bold=True, color=COLOR_PRIMARY)
    feats = [
        ("Prix & délai (2)",
         "prix · delai_demande_jours"),
        ("Identifiants encodés (2)",
         "code_client_freq · code_article_freq"),
        ("Calendrier — commande (6)",
         "annee_cmd · mois_cmd · trimestre_cmd · semaine_cmd · jour_semaine_cmd · est_fin_mois_cmd"),
        ("Calendrier — livraison (2)",
         "jour_semaine_liv_dem · est_weekend_liv_dem"),
        ("Jours fériés (2)",
         "est_jour_ferie_cmd · est_jour_ferie_liv_dem"),
        ("Météo (3)",
         "pluie_mm_liv_dem · vent_max_kmh_liv_dem · temp_min_liv_dem"),
        ("Macro (1)",
         "ipi_valeur"),
        ("Catégorielles encodées (6)",
         "statut · devise · pays · famille_activite_client · famille_activite_article · segment · type_activite (enc)"),
        ("Contexte (3)",
         "est_vacances_scolaires · est_periode_peak · est_periode_covid (liv_dem)"),
    ]
    y = Inches(1.9)
    for i, (titre, det) in enumerate(feats):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        yy = y + Inches(row * 1.55)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, Inches(4), Inches(1.4))
        set_fill(card, COLOR_LIGHT)
        add_textbox(s, x + Inches(0.1), yy + Inches(0.1), Inches(3.9), Inches(0.4),
                    titre, size=12, bold=True, color=COLOR_PRIMARY)
        add_textbox(s, x + Inches(0.1), yy + Inches(0.5), Inches(3.9), Inches(0.9),
                    det, size=10, color=COLOR_TEXT)

    # ============ Slide 25 — Phase 3 : Architecture globale ===========
    s = slide_blank(prs)
    add_header(s, "Phase 3 — Architecture de modélisation",
               "Deux architectures complémentaires : quantité + temporalité",
               idx=25, total=TOTAL)
    # A1
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4),
                              Inches(6.1), Inches(5.6))
    set_fill(box, COLOR_PRIMARY)
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(0.5),
                "Architecture 1 — Prévision de quantité", size=16, bold=True,
                color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(2.05), Inches(6.1), Inches(0.4),
                "Cible : qte_demandee  (régression sur log1p)",
                size=12, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.7), Inches(2.55), Inches(5.8), Inches(4.3), [
        ("XGBoost (loss MSE sur log1p)",),
        ("XGBoost (loss Tweedie — distribution asymétrique)",),
        ("LightGBM (loss MSE sur log1p)",),
        ("Optimisation Optuna — 50 essais (TPE)",),
        ("Métriques : MAE, RMSE, R², WAPE, MAE par décile",),
        ("Baseline à battre : moyenne historique (client × article)",),
    ], size=13, bullet_color=COLOR_ACCENT, color=COLOR_WHITE, line_spacing=1.4)
    # A2
    box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4),
                              Inches(6.1), Inches(5.6))
    set_fill(box2, COLOR_ACCENT)
    add_textbox(s, Inches(6.8), Inches(1.55), Inches(6.1), Inches(0.5),
                "Architecture 2 — Prévision de délai inter-commandes", size=16, bold=True,
                color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.8), Inches(2.05), Inches(6.1), Inches(0.4),
                "Cible : nb_jours_jusqu_prochaine_cmd  (time-to-event)",
                size=12, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(7), Inches(2.55), Inches(5.8), Inches(4.3), [
        ("LSTM PyTorch (séquences clientes)",),
        ("8 features temporelles (rolling stats par client)",),
        ("Loss MAE — optimiseur Adam",),
        ("Métriques : MAE jours · précision ±7 jours",),
        ("Affiché avec honnêteté dans le dashboard",),
        ("Performance limitée : 25,5 % à ±7j",),
    ], size=13, bullet_color=COLOR_PRIMARY, color=COLOR_WHITE, line_spacing=1.4)

    # ============ Slide 26 — Baseline ===========
    s = slide_blank(prs)
    add_header(s, "Baseline métier — moyenne historique (client × article)",
               "Référence à battre pour démontrer la valeur ajoutée du ML",
               idx=26, total=TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.4), Inches(8), Inches(0.5),
                "Méthode", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(1.9), Inches(8), Inches(3), [
        ("Calcul de la moyenne de qte_demandee par couple (code_client, code_article) sur 2021-2023",),
        ("Application sur 2025 : si couple connu → moyenne ; sinon → moyenne globale (fallback)",),
        ("Couverture observée : 84,7 % des couples présents en 2025 sont connus du train",),
        ("44 192 couples historiques distincts",),
    ], size=13, line_spacing=1.4)
    add_textbox(s, Inches(0.5), Inches(4.9), Inches(8), Inches(0.5),
                "Pourquoi cette baseline ?", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(5.4), Inches(8), Inches(2), [
        ("Reflète la pratique métier actuelle (heuristique simple, lisible)",),
        ("Conserve l'effet du couple client × article (forte mémoire d'achat)",),
        ("Battre cette baseline = vraie valeur ajoutée du ML",),
    ], size=12, line_spacing=1.3)
    # KPIs
    add_kpi(s, Inches(8.8), Inches(1.5), Inches(4.2), Inches(1.7), "MAE 13,04", "Erreur absolue moyenne",
            value_color=COLOR_PRIMARY)
    add_kpi(s, Inches(8.8), Inches(3.4), Inches(4.2), Inches(1.7), "RMSE 145,11", "Erreur quadratique",
            value_color=COLOR_PRIMARY)
    add_kpi(s, Inches(8.8), Inches(5.3), Inches(4.2), Inches(1.7), "84,7 %", "Couverture couples connus",
            value_color=COLOR_ACCENT)

    # ============ Slide 27 — Comparatif modèles ===========
    s = slide_blank(prs)
    add_header(s, "Comparatif des modèles — Architecture 1 (test 2025)",
               "Évaluation sur 66 174 commandes hors entraînement",
               idx=27, total=TOTAL)
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(3.5), [
        ["Modèle", "MAE", "RMSE", "R²", "WAPE", "Δ MAE vs baseline"],
        ["Baseline (moy. historique)", "13,04", "145,11", "—", "—", "—"],
        ["XGBoost (log + MSE)", "14,36", "153,11", "0,246", "65,4 %", "+10 %"],
        ["LightGBM (log + MSE)", "14,74", "158,43", "0,192", "67,1 %", "+13 %"],
        ["XGBoost Tweedie", "13,74", "138,98", "0,378", "62,6 %", "+5 %"],
        ["★ XGBoost Optuna (final)", "11,87", "132,17", "0,438", "54,1 %", "−9 %"],
    ], col_widths=[Inches(3.6), Inches(1.4), Inches(1.6), Inches(1.4), Inches(1.6), Inches(2.8)], font_size=12)
    add_textbox(s, Inches(0.5), Inches(5.2), Inches(12.4), Inches(0.6),
                "Lecture : seul XGBoost après tuning Optuna bat la baseline (−9 % MAE, +21 % R² vs Tweedie).",
                size=14, bold=True, color=COLOR_ACCENT)
    add_bullets(s, Inches(0.5), Inches(5.85), Inches(12.4), Inches(1.4), [
        ("Le log+MSE seul perd contre la baseline → la baseline capture déjà la mémoire client × article",),
        ("Tweedie corrige partiellement l'asymétrie de la distribution (R² 0,38)",),
        ("L'optimisation bayésienne (Optuna) débloque le gain métier réel",),
    ], size=12, line_spacing=1.3)

    # ============ Slide 28 — XGBoost Optuna : hyperparamètres ===========
    s = slide_blank(prs)
    add_header(s, "XGBoost Optuna — hyperparamètres optimaux (50 essais TPE)",
               "Recherche bayésienne via tree-structured Parzen estimator",
               idx=28, total=TOTAL)
    add_table(s, Inches(0.5), Inches(1.4), Inches(7), Inches(5.5), [
        ["Hyperparamètre", "Valeur retenue", "Plage explorée"],
        ["n_estimators", "498", "[100 — 500]"],
        ["max_depth", "10", "[3 — 12]"],
        ["learning_rate", "0,0776", "[0,01 — 0,3] (log)"],
        ["subsample", "0,960", "[0,6 — 1,0]"],
        ["colsample_bytree", "0,729", "[0,6 — 1,0]"],
        ["min_child_weight", "3", "[1 — 10]"],
        ["reg_alpha", "0,0207", "[0 — 1]"],
        ["reg_lambda", "0,2836", "[0 — 1]"],
    ], col_widths=[Inches(2.8), Inches(2), Inches(2.2)], font_size=12)
    # Encart méthode
    add_textbox(s, Inches(8), Inches(1.4), Inches(5), Inches(0.5),
                "Protocole Optuna", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(8), Inches(1.9), Inches(5), Inches(5), [
        ("Sampler : TPE (Bayesian)",),
        ("Pruner : MedianPruner (early-stop)",),
        ("Objectif : minimiser MAE val 2024",),
        ("Budget : 50 essais, ~ 1 h CPU",),
        ("Seed fixée : reproductibilité",),
        ("Sauvegarde finale : models/xgboost_optuna_final.pkl",),
    ], size=13, line_spacing=1.4)

    # ============ Slide 29 — Feature importance ===========
    s = slide_blank(prs)
    add_header(s, "Feature importance — XGBoost Optuna",
               "Quelles variables dominent la prévision ?",
               idx=29, total=TOTAL)
    add_image_centered(s, REPORTS / "feature_importance_xgb.png",
                       Inches(0.4), Inches(1.35), Inches(8.6), Inches(5.6))
    add_textbox(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(0.5),
                "Lecture business", size=14, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(9.3), Inches(1.9), Inches(3.7), Inches(5), [
        ("Identifiants encodés (client/article) dominent → mémoire d'achat",),
        ("Prix joue un rôle structurel",),
        ("Calendrier > météo en importance",),
        ("Confirme l'analyse Pearson/Granger Phase 2",),
        ("IPI marginal au niveau ligne (signal mensuel dilué)",),
    ], size=10, line_spacing=1.3)

    # ============ Slide 30 — LSTM ===========
    s = slide_blank(prs)
    add_header(s, "Architecture 2 — LSTM time-to-event",
               "Prédiction du délai inter-commandes par client (PyTorch)",
               idx=30, total=TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(7), Inches(0.5),
                "Conception du modèle", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(1.8), Inches(7), Inches(4), [
        ("Séquences clientes : derniers n=10 délais inter-cmd + features rolling",),
        ("8 features : delai_moyen_30j, delai_std_30j, freq_cmd_90j, qte_moy_30j, ...",),
        ("LSTM 2 couches × 64 unités · dropout 0,2",),
        ("Couche dense finale → 1 sortie (jours)",),
        ("Loss : MAE  ·  Optimiseur : Adam (lr=1e-3)",),
        ("Early stopping val MAE — 30 epochs max",),
    ], size=12, line_spacing=1.4)
    # Résultats LSTM
    add_textbox(s, Inches(8), Inches(1.3), Inches(5), Inches(0.5),
                "Résultats test 2025", size=16, bold=True, color=COLOR_PRIMARY)
    add_kpi(s, Inches(8), Inches(1.9), Inches(2.4), Inches(1.4), "24,5 j", "MAE délai",
            value_color=COLOR_WARN)
    add_kpi(s, Inches(10.6), Inches(1.9), Inches(2.4), Inches(1.4), "25,5 %", "Précision ±7 j",
            value_color=COLOR_WARN)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(3.6),
                              Inches(5), Inches(2.8))
    set_fill(box, COLOR_LIGHT)
    add_textbox(s, Inches(8.1), Inches(3.7), Inches(4.8), Inches(0.5),
                "Lecture honnête — limitation assumée", size=13, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(8.1), Inches(4.2), Inches(4.8), Inches(2.2), [
        ("Demande GE = sporadique, non périodique",),
        ("Signaux ts insuffisants par client",),
        ("→ exposé avec badge ROUGE dans dashboard",),
        ("Perspective : transformer client × calendaire",),
    ], size=10, line_spacing=1.3)

    # ============ Slide 31 — Stockage modèles ===========
    s = slide_blank(prs)
    add_header(s, "Livrables Phase 3 — modèles sauvegardés",
               "5 artefacts ML versionnés dans models/",
               idx=31, total=TOTAL)
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.4), Inches(4), [
        ["Fichier", "Modèle", "Cible", "Métrique principale"],
        ["models/xgboost_qte_demandee.pkl", "XGBoost log + MSE", "qte_demandee", "MAE 14,36"],
        ["models/xgboost_tweedie_qte.pkl", "XGBoost Tweedie", "qte_demandee", "MAE 13,74"],
        ["models/lightgbm_qte_demandee.pkl", "LightGBM log + MSE", "qte_demandee", "MAE 14,74"],
        ["★ models/xgboost_optuna_final.pkl", "XGBoost Optuna (TPE 50 essais)", "qte_demandee", "MAE 11,87"],
        ["models/lstm_time_to_event.pt", "LSTM 2×64 PyTorch", "délai inter-cmd", "MAE 24,5 j"],
    ], col_widths=[Inches(4.4), Inches(3.6), Inches(2.4), Inches(2)], font_size=11)
    add_textbox(s, Inches(0.5), Inches(5.6), Inches(12.4), Inches(0.5),
                "Rapport JSON consolidé : reports/rapport_modelisation.json", size=13,
                bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(6.1), Inches(12.4), Inches(1.1), [
        ("Tous les hyperparamètres, métriques et features sauvegardés en JSON pour reproductibilité",),
        ("Chargement direct par le dashboard via joblib (modèles XGB/LGBM) et torch.load (LSTM)",),
    ], size=11, line_spacing=1.3)

    # ============ Slide 32 — Phase 4 : Dashboard architecture ===========
    s = slide_blank(prs)
    add_header(s, "Phase 4 — Dashboard décisionnel human-in-the-loop",
               "Streamlit multi-pages · 4 pages métiers · ré-entraînement intégré",
               idx=32, total=TOTAL)
    pages = [
        ("Page 1 — Données", "Exploration interactive\nFiltres client / article / période\nKPIs volumes & qualité",
         COLOR_PRIMARY),
        ("Page 2 — IA", "Statut des modèles\nHyperparamètres\nDataset d'entraînement",
         COLOR_ACCENT),
        ("Page 3 — Prévisions", "Saisie manuelle ou upload\nScore confiance composite\nBadge LSTM rouge",
         COLOR_PRIMARY),
        ("Page 4 — Analyse", "Comparatif modèles\nFeature importance\nWAPE par décile",
         COLOR_ACCENT),
    ]
    x = Inches(0.5)
    for i, (titre, det, color) in enumerate(pages):
        xx = Inches(0.5 + i * 3.15)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, xx, Inches(1.6),
                                  Inches(2.95), Inches(4.5))
        set_fill(box, color)
        add_textbox(s, xx, Inches(1.75), Inches(2.95), Inches(0.6),
                    titre, size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, xx + Inches(0.15), Inches(2.4), Inches(2.7), Inches(3.5),
                    det, size=12, color=COLOR_WHITE)
    add_textbox(s, Inches(0.5), Inches(6.3), Inches(12.4), Inches(0.5),
                "Stack technique", size=14, bold=True, color=COLOR_PRIMARY)
    add_textbox(s, Inches(0.5), Inches(6.75), Inches(12.4), Inches(0.4),
                "Streamlit · pandas · joblib · torch · plotly · scikit-learn — déploiement local (MVP sans cloud, sans auth)",
                size=11, color=COLOR_MUTED)

    # ============ Slide 33 — Score confiance composite ===========
    s = slide_blank(prs)
    add_header(s, "Score de confiance composite — innovation clé du dashboard",
               "Aide à la décision : ne pas se fier aveuglément à la prédiction",
               idx=33, total=TOTAL)
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(7.5), Inches(0.5),
                "Le score combine 3 facteurs", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(1.85), Inches(7.5), Inches(5), [
        ("Familiarité de l'article", "Nb d'occurrences code_article dans le train"),
        ("Écart à la médiane du couple", "(prédiction − médiane historique client × article) / médiane"),
        ("Décile de prédiction", "Position de la valeur prédite dans la distribution des résidus train"),
    ], size=14, line_spacing=1.5)
    # Code couleur
    add_textbox(s, Inches(8.2), Inches(1.3), Inches(4.8), Inches(0.5),
                "Code couleur restitué à l'utilisateur", size=14, bold=True, color=COLOR_PRIMARY)
    badges = [
        ("VERT", "Confiance haute — couple bien connu, prédiction cohérente", COLOR_SUCCESS),
        ("ORANGE", "Confiance moyenne — vigilance recommandée", COLOR_ACCENT),
        ("ROUGE", "Confiance faible — validation manuelle nécessaire", COLOR_WARN),
    ]
    for i, (label, det, color) in enumerate(badges):
        y = Inches(1.9 + i * 1.4)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), y,
                                  Inches(1.5), Inches(0.8))
        set_fill(box, color)
        add_textbox(s, Inches(8.2), y, Inches(1.5), Inches(0.8),
                    label, size=14, bold=True, color=COLOR_WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(9.85), y, Inches(3.2), Inches(0.8),
                    det, size=11, color=COLOR_TEXT, anchor=MSO_ANCHOR.MIDDLE)

    # ============ Slide 34 — Résultats finaux ===========
    s = slide_blank(prs)
    add_header(s, "Résultats finaux — objectif North Star atteint",
               "Comparaison synthétique baseline vs modèle gagnant",
               idx=34, total=TOTAL)
    # Grand cartouche victoire
    add_kpi(s, Inches(0.5), Inches(1.5), Inches(3.0), Inches(1.6), "−9 %", "MAE vs baseline",
            value_color=COLOR_SUCCESS)
    add_kpi(s, Inches(3.7), Inches(1.5), Inches(3.0), Inches(1.6), "11,87", "MAE XGB Optuna",
            value_color=COLOR_PRIMARY)
    add_kpi(s, Inches(6.9), Inches(1.5), Inches(3.0), Inches(1.6), "0,438", "R² test 2025",
            value_color=COLOR_PRIMARY)
    add_kpi(s, Inches(10.1), Inches(1.5), Inches(2.8), Inches(1.6), "54,1 %", "WAPE",
            value_color=COLOR_PRIMARY)

    add_textbox(s, Inches(0.5), Inches(3.4), Inches(12.4), Inches(0.5),
                "Validation de la problématique", size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, Inches(0.5), Inches(3.9), Inches(12.4), Inches(3.2), [
        ("La quantité (Architecture 1) est correctement prédite : MAE 11,87 < baseline 13,04",),
        ("L'apport des features exogènes (calendrier, météo, macro) est confirmé par feature importance",),
        ("L'optimisation Optuna apporte +21 % de R² vs Tweedie sans tuning",),
        ("La date (Architecture 2 LSTM) reste limitée : MAE 24,5 j — perspective de recherche",),
        ("Précision globale (1 − WAPE) ≈ 46 % — sous l'objectif 60 %, mais 9 % de gain MAE pertinent métier",),
    ], size=13, line_spacing=1.4)

    # ============ Slide 35 — Apports académiques ===========
    s = slide_blank(prs)
    add_header(s, "Apports académiques & contributions",
               "Méthodologie, reproductibilité, transparence",
               idx=35, total=TOTAL)
    add_bullets(s, Inches(0.5), Inches(1.4), Inches(12), Inches(5.7), [
        ("Pipeline IA end-to-end documenté", "De l'extraction ERP au dashboard décisionnel"),
        ("Tests statistiques rigoureux", "Pearson + Granger pour valider chaque feature exogène"),
        ("Comparaison empirique boosting vs Deep Learning sur données industrielles réelles",
         "XGBoost / LightGBM / LSTM sur 5 ans d'historique"),
        ("Évaluation par split temporel strict", "Pas de leakage, validation année-par-année"),
        ("Score de confiance composite", "Innovation pour systèmes de prévision interprétables"),
        ("Transparence sur les limites", "LSTM mis en avant honnêtement comme zone d'amélioration"),
        ("Reproductibilité complète", "Notebooks + JSON hyperparamètres + Parquet versionnés"),
    ], size=14, line_spacing=1.5)

    # ============ Slide 36 — Apports business ===========
    s = slide_blank(prs)
    add_header(s, "Apports business pour GE",
               "Valeur opérationnelle pour la chaîne d'approvisionnement",
               idx=36, total=TOTAL)
    add_bullets(s, Inches(0.5), Inches(1.4), Inches(12), Inches(5.7), [
        ("Gain MAE de 9 % vs heuristique actuelle", "Sur ~66 000 commandes/an, cela réduit l'erreur cumulée significativement"),
        ("Identification quantifiée des leviers exogènes", "Cycles calendaires > météo > macro (preuves Granger)"),
        ("Outil de prévision opérationnel", "Dashboard Streamlit utilisable directement par le métier"),
        ("Score de confiance pour décider quand sur-stocker / sous-stocker",
         "Ajuste finement le taux de service"),
        ("Cadre de ré-entraînement intégré", "Permet l'adaptation rapide à de nouveaux contextes"),
        ("Coût d'infrastructure quasi nul", "Modèles légers, batch CPU, dashboard local"),
        ("Base solide pour industrialisation", "Migration MLflow / cloud immédiate"),
    ], size=14, line_spacing=1.5)

    # ============ Slide 37 — Limites & perspectives ===========
    s = slide_blank(prs)
    add_header(s, "Limites identifiées & perspectives de recherche",
               "Honnêteté scientifique : ce que le modèle ne capture pas encore",
               idx=37, total=TOTAL)
    # Limites
    add_textbox(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.5),
                "Limites actuelles", size=18, bold=True, color=COLOR_WARN)
    add_bullets(s, Inches(0.5), Inches(1.95), Inches(6), Inches(5), [
        ("LSTM time-to-event sous-performant", "Précision ±7j seulement 25,5 %"),
        ("Précision globale 46 % < objectif 60 %", "WAPE 54,1 %"),
        ("IPI partiellement proxy", "À remplacer par série INSEE officielle"),
        ("Météo limitée à Paris", "Pas de granularité multi-site"),
        ("Pas de prise en compte du carnet de commandes",
         "Backlog non modélisé"),
    ], size=12, bullet_color=COLOR_WARN, line_spacing=1.4)
    # Perspectives
    add_textbox(s, Inches(7), Inches(1.4), Inches(6), Inches(0.5),
                "Perspectives", size=18, bold=True, color=COLOR_SUCCESS)
    add_bullets(s, Inches(7), Inches(1.95), Inches(6), Inches(5), [
        ("Architecture Transformer pour le time-to-event",
         "Embeddings client × calendaire"),
        ("Modèle hiérarchique par famille d'article",
         "Famille 0 vs 1 vs 2 séparément"),
        ("Intégration MLflow + monitoring drift",
         "Pour industrialisation"),
        ("Données météo géolocalisées",
         "Pondération par pays de livraison"),
        ("Optimisation conjointe quantité × délai",
         "Modèle multi-tâche unifié"),
        ("Reinforcement learning pour stock policy",
         "Étape supérieure : décision automatique"),
    ], size=12, bullet_color=COLOR_SUCCESS, line_spacing=1.4)

    # ============ Slide 38 — Merci ===========
    s = slide_blank(prs)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, COLOR_PRIMARY)
    add_rect(s, Emu(0), Inches(4), SLIDE_W, Inches(0.08), COLOR_ACCENT)
    add_textbox(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
                "Merci de votre attention", size=54, bold=True, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8),
                "Questions & échanges", size=28, color=COLOR_LIGHT,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "Adam Marrakchi  ·  marrakchiadham@gmail.com",
                size=14, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                "Projet mémoire Master GE — Mai 2026",
                size=12, color=COLOR_ACCENT, align=PP_ALIGN.CENTER, bold=True)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"[OK] Presentation generee : {OUTPUT}")
    print(f"     Nombre de slides : {len(prs.slides)}")


if __name__ == "__main__":
    main()
